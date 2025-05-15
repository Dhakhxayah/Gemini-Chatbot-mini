import streamlit as st
import pandas as pd
import numpy as np
import google.generativeai as genai
from sklearn.metrics.pairwise import cosine_similarity
from sentence_transformers import SentenceTransformer
import json
import pptx
import PyPDF2

# ✅ Page config
st.set_page_config(page_title="Chatbot Mini", page_icon="🏙️")

# ✅ API Key config
genai.configure(api_key="Your api key")
model = genai.GenerativeModel('gemini-2.0-flash')

# ✅ Embedding model
embedding_model = SentenceTransformer('all-MiniLM-L6-v2')

# 🔁 Session State
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "file_text" not in st.session_state:
    st.session_state.file_text = ""
if "csv_embeddings" not in st.session_state:
    st.session_state.csv_embeddings = None
    st.session_state.csv_df = None

# 🔲 Clear Chat Button
if st.button("🧹 Clear Chat"):
    st.session_state.chat_history.clear()

# 📄 File Reader
def read_file(file):
    try:
        if file.type == "text/plain":
            return file.read().decode("utf-8")
        elif file.type == "application/json":
            return json.dumps(json.load(file), indent=2)
        elif file.type == "application/pdf":
            reader = PyPDF2.PdfReader(file)
            return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])
        elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
            prs = pptx.Presentation(file)
            return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif file.type == "text/csv":
            df = pd.read_csv(file)
            combined = df.astype(str).agg(' | '.join, axis=1)
            embeddings = [embedding_model.encode(text) for text in combined]
            st.session_state.csv_df = df
            st.session_state.csv_embeddings = embeddings
            return "\n".join(combined)
        else:
            return "Unsupported file."
    except Exception as e:
        return f"❌ Error reading file: {e}"

# 🔍 Context Fetch
def fetch_context(question):
    if st.session_state.csv_embeddings:
        question_vec = embedding_model.encode(question)
        similarities = cosine_similarity([question_vec], st.session_state.csv_embeddings)[0]
        top_indices = np.argsort(similarities)[::-1][:5]
        top_rows = st.session_state.csv_df.iloc[top_indices]
        return "\n\n".join(top_rows.astype(str).agg(' | '.join, axis=1))
    else:
        return st.session_state.file_text

# 🧠 Gemini Bot
def ask_bot(question):
    context = fetch_context(question)
    prompt = f"Use the following content to answer:\n\n{context}\n\nQuestion: {question}" if context else question
    try:
        response = model.generate_content(prompt)
        return response.text if response else "⚠️ No response from Gemini API"
    except Exception as e:
        return f"⚠️ Error: {e}"

# 🧾 Display Chat History
st.markdown("### 💬 Chat")
for chat in st.session_state.chat_history:
    st.markdown(f"**🧑 You**: {chat['user']}")
    st.markdown(f"**🤖 Bot**: {chat['bot']}")

# 🧠 Chat Input Form with Upload at Bottom
st.markdown("---")
with st.form("chat_form", clear_on_submit=True):
    cols = st.columns([0.85, 0.15])
    with cols[0]:
        user_input = st.text_input("💬 Ask me anything:")
    with cols[1]:
        uploaded_file = st.file_uploader("📁", type=['txt', 'csv', 'pdf', 'json', 'pptx'], label_visibility="collapsed")
    submitted = st.form_submit_button("Send")

# When Send button is clicked
if submitted:
    if uploaded_file:
        st.session_state.file_text = read_file(uploaded_file)
        st.success("✅ File uploaded and processed!")

    if user_input:
        with st.spinner("🤖 Thinking..."):
            answer = ask_bot(user_input)
            st.session_state.chat_history.append({"user": user_input, "bot": answer})
        st.experimental_rerun()
